#!/usr/bin/env python3
"""Build an Evidura-branded PowerPoint deck from a JSON deck spec.

    python3 scripts/build-evidura-deck.py <spec.json> -o <out.pptx>

The split this script exists to enforce: **an agent writes the content, this
script owns the brand.** The spec carries words and numbers only. Every colour,
every rule, the mark, the type scale and the slide furniture come from here and
from `brand/evidura/tokens.css`, so a brand change lands in one place and no
agent is ever asked to remember a hex value.

Spec shape (see docs/dfva-v4-socialisation-harness.md for the full contract):

    {
      "program": {"code": "mc-mgmthre", "name": "Master of ..."},
      "meta": {"instrument": "4.1-draft", "assessmentDate": "2026-08-14"},
      "slides": [{"layout": "title", ...}, ...]
    }

Layouts: title · statement · bullets · two-col · scores · findings · table ·
questions · closing. An unknown layout is a hard error — a silently skipped
slide is worse than a failed build.
"""

from __future__ import annotations

import argparse
import json
import math
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parent.parent
TOKENS = REPO / "brand" / "evidura" / "tokens.css"

# Fallbacks are the published token values. They exist so a missing tokens.css
# fails loudly (below) rather than producing an unbranded deck that looks fine.
FALLBACK = {
    "ink": "0F2236",
    "signal": "E9A23B",
    "paper": "F4F2EC",
    "slate": "5C7088",
    "ink-900": "0A1828",
    "ink-700": "1B3554",
    "ink-100": "DDE3EB",
    "signal-700": "B97E26",
    "paper-200": "E7E3D8",
    "band-resilient": "1F9D6B",
    "band-moderate": "E9A23B",
    "band-high": "E0742F",
    "band-critical": "C7433A",
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.85)
# Average glyph advance as a fraction of the point size, used by every text
# measurement below. Deliberately pessimistic: PowerPoint substitutes a wider
# face when Inter is absent, and an optimistic constant overflows on the
# machine that opens the deck rather than the one that built it.
GLYPH = 0.55
# Rendered line height exceeds size x line-spacing by the font's own ascent and
# descent. Without this the estimate clears the box and the text still overruns.
LEADING = 1.2
CONTENT_W = SLIDE_W - 2 * MARGIN


def load_tokens(path: Path) -> dict[str, RGBColor]:
    """Parse --evidura-* hex custom properties out of tokens.css."""
    values = dict(FALLBACK)
    if path.exists():
        css = path.read_text(encoding="utf-8")
        for name, hexval in re.findall(r"--evidura-([a-z0-9-]+):\s*#([0-9A-Fa-f]{6})", css):
            values[name] = hexval.upper()
    else:
        print(f"warning: {path} not found — falling back to embedded token values", file=sys.stderr)
    return {k: RGBColor.from_string(v) for k, v in values.items()}


class Deck:
    def __init__(self, spec: dict, font: str, tokens: dict[str, RGBColor]) -> None:
        self.spec = spec
        self.font = font
        self.c = tokens
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.program = spec.get("program", {})
        self.meta = spec.get("meta", {})
        self.n = 0
        self.overflows: list[str] = []

    # ---------- primitives ----------

    def _fit(self, text, width, height, start, *, min_pt=9.0, spacing=1.2, where=""):
        """Largest point size at or below `start` whose wrapped text fits the box.

        PowerPoint's own autofit only applies on edit, so a deck built headlessly
        overflows silently — which is how four slides shipped with text running
        out the bottom of their cards on the first build of this script. The
        estimate is deliberately conservative (see GLYPH); when even
        `min_pt` will not fit, the slide is recorded and reported at the end so
        the author trims the copy rather than the reader losing a line.
        """
        w_in = width / 914400
        h_in = height / 914400
        paras = text if isinstance(text, list) else [text]
        pt = start
        while pt >= min_pt:
            cpl = max(int((w_in * 72) / (GLYPH * pt)), 1)
            lines = sum(max(1, math.ceil(len(p) / cpl)) for p in paras)
            if lines * pt * spacing * LEADING / 72 <= h_in:
                return pt
            pt -= 0.5
        self.overflows.append(f"slide {self.n + 1} ({where}): copy too long — trim it")
        return min_pt

    def _slide(self, bg: RGBColor):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bgshape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        bgshape.fill.solid()
        bgshape.fill.fore_color.rgb = bg
        bgshape.line.fill.background()
        bgshape.shadow.inherit = False
        return s

    def _text(
        self,
        slide,
        text,
        left,
        top,
        width,
        height,
        *,
        size=18,
        color=None,
        bold=False,
        italic=False,
        align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,
        spacing=1.15,
        font=None,
        caps=False,
    ):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        lines = text if isinstance(text, list) else [text]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.line_spacing = spacing
            run = p.add_run()
            run.text = line.upper() if caps else line
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.name = font or self.font
            run.font.color.rgb = color if color is not None else self.c["ink"]
        return box

    def _rule(self, slide, left, top, width, color, thickness=Pt(3)):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Emu(int(thickness)))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        bar.shadow.inherit = False
        return bar

    def _pill(self, slide, left, top, width, height, fill, line=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.5
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if line is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = line
            shape.line.width = Pt(1)
        shape.shadow.inherit = False
        return shape

    def _card(self, slide, left, top, width, height, fill, line=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.06
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if line is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = line
            shape.line.width = Pt(1)
        shape.shadow.inherit = False
        return shape

    def _mark(self, slide, left, top, scale=1.0, on_dark=False):
        """The Strata-E mark, drawn as native shapes from the SVG geometry.

        Source: brand/evidura/evidura-mark.svg — three rounded bars on a
        200-unit canvas, the top bar in Signal and the lower two in Ink
        (Paper when reversed). Redrawn rather than embedded because PowerPoint
        cannot render SVG and a rasterised mark would go soft on a projector.
        """
        unit = Inches(0.09) * scale  # one bar height (18/200 of the canvas)
        ink = self.c["paper"] if on_dark else self.c["ink"]
        widths = [(120, self.c["signal"]), (84, ink), (120, ink)]
        for i, (w, colour) in enumerate(widths):
            bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left,
                int(top + i * unit * 1.83),
                int(unit * (w / 18)),
                int(unit),
            )
            bar.adjustments[0] = 0.5
            bar.fill.solid()
            bar.fill.fore_color.rgb = colour
            bar.line.fill.background()
            bar.shadow.inherit = False

    def _chrome(self, slide, on_dark=False):
        """Mark, wordmark, program footer and slide number on every slide."""
        self.n += 1
        muted = self.c["slate"] if not on_dark else RGBColor.from_string("8FA3B8")
        self._mark(slide, MARGIN, SLIDE_H - Inches(0.78), scale=0.42, on_dark=on_dark)
        self._text(
            slide,
            "EVIDURA",
            MARGIN + Inches(0.52),
            SLIDE_H - Inches(0.75),
            Inches(3),
            Inches(0.3),
            size=10,
            color=muted,
            bold=True,
        )
        code = self.program.get("code", "").upper()
        self._text(
            slide,
            f"{code} · durability assessment · confidential draft",
            MARGIN + Inches(1.6),
            SLIDE_H - Inches(0.75),
            CONTENT_W - Inches(2.2),
            Inches(0.3),
            size=10,
            color=muted,
        )
        self._text(
            slide,
            str(self.n),
            SLIDE_W - MARGIN - Inches(0.6),
            SLIDE_H - Inches(0.75),
            Inches(0.6),
            Inches(0.3),
            size=10,
            color=muted,
            align=PP_ALIGN.RIGHT,
        )

    @staticmethod
    def _lines(text, width, size):
        cpl = max(int(((width / 914400) * 72) / (GLYPH * size)), 1)
        return max(1, math.ceil(len(text) / cpl))

    def _heading(self, slide, title, subtitle=None, on_dark=False):
        """Title, optional subtitle, signal rule. Returns the y where content starts.

        Both blocks are measured rather than assumed: a two-line subtitle used to
        run under the rule, and a long title used to run into the subtitle.
        """
        fg = self.c["paper"] if on_dark else self.c["ink"]
        muted = RGBColor.from_string("8FA3B8") if on_dark else self.c["slate"]
        size = 34 if len(title) < 58 else 28
        title_h = int(Inches(0.52) * self._lines(title, CONTENT_W, size))
        top = Inches(0.72)
        self._text(slide, title, MARGIN, top, CONTENT_W, title_h, size=size, bold=True, color=fg)
        top = int(top + title_h + Inches(0.18))
        if subtitle:
            sub_h = int(Inches(0.3) * self._lines(subtitle, CONTENT_W, 15))
            self._text(slide, subtitle, MARGIN, top, CONTENT_W, sub_h, size=15, color=muted)
            top = int(top + sub_h + Inches(0.16))
        self._rule(slide, MARGIN, top, Inches(1.5), self.c["signal"])
        return int(top + Inches(0.42))

    def _band_colour(self, score, maximum):
        if maximum <= 0:
            return self.c["slate"]
        ratio = score / maximum
        if ratio >= 0.99:
            return self.c["band-resilient"]
        if ratio >= 0.66:
            return self.c["band-moderate"]
        if ratio >= 0.34:
            return self.c["band-high"]
        return self.c["band-critical"]

    # ---------- layouts ----------

    def title(self, s):
        slide = self._slide(self.c["ink"])
        self._mark(slide, MARGIN, Inches(1.15), scale=1.5, on_dark=True)
        self._text(
            slide, "EVIDURA", MARGIN + Inches(1.85), Inches(1.32), Inches(6), Inches(0.6),
            size=26, bold=True, color=self.c["paper"],
        )
        self._text(
            slide, s.get("eyebrow", "Durability assessment"), MARGIN + Inches(1.9), Inches(1.92),
            Inches(8), Inches(0.4), size=13, color=self.c["signal"], caps=True,
        )
        self._rule(slide, MARGIN, Inches(2.95), Inches(2.2), self.c["signal"], thickness=Pt(4))
        self._text(
            slide, s["title"], MARGIN, Inches(3.35), CONTENT_W - Inches(1.5), Inches(1.6),
            size=42 if len(s["title"]) < 52 else 34, bold=True, color=self.c["paper"], spacing=1.05,
        )
        if s.get("subtitle"):
            self._text(
                slide, s["subtitle"], MARGIN, Inches(5.05), CONTENT_W - Inches(2.4), Inches(0.7),
                size=17, color=RGBColor.from_string("8FA3B8"),
            )
        meta = s.get("meta", [])
        if meta:
            self._text(
                slide, "   ·   ".join(meta), MARGIN, Inches(5.85), CONTENT_W, Inches(0.4),
                size=12, color=RGBColor.from_string("8FA3B8"),
            )
        if s.get("stamp"):
            pill = self._pill(
                slide, SLIDE_W - MARGIN - Inches(3.6), Inches(1.25), Inches(3.6), Inches(0.42),
                self.c["ink-700"], line=self.c["signal"],
            )
            tf = pill.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = s["stamp"].upper()
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.name = self.font
            run.font.color.rgb = self.c["signal"]
        self._chrome(slide, on_dark=True)

    def statement(self, s):
        slide = self._slide(self.c["ink"])
        if s.get("kicker"):
            self._text(
                slide, s["kicker"], MARGIN, Inches(1.5), CONTENT_W, Inches(0.4),
                size=13, color=self.c["signal"], bold=True, caps=True,
            )
        self._rule(slide, MARGIN, Inches(2.05), Inches(1.5), self.c["signal"])
        body = s["statement"]
        body_w = CONTENT_W - Inches(1.2)
        self._text(
            slide, body, MARGIN, Inches(2.5), body_w, Inches(3.15),
            size=self._fit(body, body_w, Inches(3.15), 30, min_pt=17, spacing=1.22, where="statement"),
            color=self.c["paper"], spacing=1.22,
        )
        if s.get("attribution"):
            self._text(
                slide, s["attribution"], MARGIN, Inches(5.9), CONTENT_W, Inches(0.5),
                size=13, color=RGBColor.from_string("8FA3B8"), italic=True,
            )
        self._chrome(slide, on_dark=True)

    def bullets(self, s):
        slide = self._slide(self.c["paper"])
        top = self._heading(slide, s["title"], s.get("subtitle"))
        items = s.get("bullets", [])
        gap = Inches(0.16)
        avail = SLIDE_H - top - Inches(1.15) - (Inches(0.42) if s.get("note") else 0)
        h = int((avail - gap * max(len(items) - 1, 0)) / max(len(items), 1))
        for i, item in enumerate(items):
            y = int(top + i * (h + gap))
            self._rule(slide, MARGIN, y + Inches(0.12), Inches(0.22), self.c["signal"], thickness=Pt(3))
            lead = item.get("lead")
            text_left = MARGIN + Inches(0.45)
            body_w = CONTENT_W - Inches(0.45)
            if lead:
                self._text(slide, lead, text_left, y, body_w, Inches(0.34), size=16, bold=True)
                body_h = h - Inches(0.38)
                self._text(
                    slide, item.get("text", ""), text_left, y + Inches(0.38), body_w, body_h,
                    size=self._fit(item.get("text", ""), body_w, body_h, 13.5, where="bullets"),
                    color=self.c["slate"],
                )
            else:
                self._text(
                    slide, item.get("text", ""), text_left, y, body_w, h,
                    size=self._fit(item.get("text", ""), body_w, h, 15, where="bullets"),
                )
        if s.get("note"):
            self._text(
                slide, s["note"], MARGIN, SLIDE_H - Inches(1.22), CONTENT_W, Inches(0.4),
                size=12, color=self.c["slate"], italic=True,
            )
        self._chrome(slide)

    def two_col(self, s):
        slide = self._slide(self.c["paper"])
        top = self._heading(slide, s["title"], s.get("subtitle"))
        col_w = int((CONTENT_W - Inches(0.5)) / 2)
        height = SLIDE_H - top - Inches(1.2)
        for i, side in enumerate(("left", "right")):
            col = s.get(side, {})
            x = int(MARGIN + i * (col_w + Inches(0.5)))
            accent = self.c["signal"] if i == 0 else self.c["slate"]
            self._card(slide, x, top, col_w, height, self.c["paper-200"])
            self._rule(slide, x, top, col_w, accent, thickness=Pt(4))
            self._text(
                slide, col.get("heading", ""), x + Inches(0.32), top + Inches(0.32),
                col_w - Inches(0.64), Inches(0.4), size=17, bold=True,
            )
            items = [f"—  {t}" for t in col.get("items", [])]
            body_w = col_w - Inches(0.64)
            body_h = height - Inches(1.2)
            self._text(
                slide, items, x + Inches(0.32), top + Inches(0.92), body_w, body_h,
                size=self._fit(items, body_w, body_h, 13.5, spacing=1.45, where="two-col"),
                color=self.c["ink"], spacing=1.45,
            )
        self._chrome(slide)

    def scores(self, s):
        slide = self._slide(self.c["paper"])
        top = self._heading(slide, s["title"], s.get("subtitle"))

        floor = SLIDE_H - Inches(1.05)  # everything must clear the footer chrome

        # Left rail: exposure, then each sub-score in its own card, then the gates.
        # The rail is the one place a combined figure could appear by accident, so
        # the sub-scores are drawn as separate cards and never share a container.
        rail_w = Inches(3.5)
        exp = s.get("exposure")
        y = top
        if exp:
            self._card(slide, MARGIN, y, rail_w, Inches(1.5), self.c["ink"])
            self._text(
                slide, exp.get("label", "Destination AI exposure"), MARGIN + Inches(0.28),
                y + Inches(0.18), rail_w - Inches(0.56), Inches(0.28), size=10.5,
                color=RGBColor.from_string("8FA3B8"), caps=True, bold=True,
            )
            self._text(
                slide, str(exp.get("value", "")), MARGIN + Inches(0.28), y + Inches(0.48),
                rail_w - Inches(0.56), Inches(0.62), size=38, bold=True, color=self.c["signal"],
            )
            self._text(
                slide, exp.get("note", ""), MARGIN + Inches(0.28), y + Inches(1.12),
                rail_w - Inches(0.4), Inches(0.3), size=10, color=RGBColor.from_string("8FA3B8"),
            )
            y = int(y + Inches(1.66))
        for sub in s.get("subscales", []):
            self._card(slide, MARGIN, y, rail_w, Inches(0.86), self.c["paper-200"])
            self._rule(slide, MARGIN, y, rail_w, self.c["ink"], thickness=Pt(3))
            self._text(
                slide, sub.get("label", ""), MARGIN + Inches(0.28), y + Inches(0.16),
                rail_w - Inches(0.56), Inches(0.26), size=11.5, color=self.c["slate"], bold=True,
            )
            self._text(
                slide, str(sub.get("value", "")), MARGIN + Inches(0.28), y + Inches(0.42),
                rail_w - Inches(0.56), Inches(0.38), size=22, bold=True,
            )
            y = int(y + Inches(1.0))

        gates = s.get("gates", [])
        if gates:
            gw = int((rail_w - Inches(0.2) * (len(gates) - 1)) / len(gates))
            for i, gate in enumerate(gates):
                gx = int(MARGIN + i * (gw + Inches(0.2)))
                passed = str(gate.get("result", "")).upper() == "PASS"
                colour = self.c["band-resilient"] if passed else self.c["band-critical"]
                pill = self._pill(slide, gx, y, gw, Inches(0.42), self.c["paper-200"], line=colour)
                p = pill.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = f"{gate.get('id', '')}  {gate.get('result', '')}"
                run.font.size = Pt(12)
                run.font.bold = True
                run.font.name = self.font
                run.font.color.rgb = colour

        # Right grid: one chip per item, coloured by level against the band scale.
        grid_x = int(MARGIN + rail_w + Inches(0.55))
        grid_w = int(SLIDE_W - MARGIN - grid_x)
        items = s.get("items", [])
        cols = 2
        cw = int((grid_w - Inches(0.3)) / cols)
        rows = (len(items) + cols - 1) // cols
        notes = [n for n in (s.get("railNote"), s.get("note")) if n]
        note_h = Inches(0.5) if notes else 0
        ch = min(Inches(0.78), int((floor - top - note_h - Inches(0.18) * max(rows - 1, 0)) / max(rows, 1)))
        for i, item in enumerate(items):
            r, c = divmod(i, cols)
            x = int(grid_x + c * (cw + Inches(0.3)))
            yy = int(top + r * (ch + Inches(0.18)))
            maximum = item.get("max", 3)
            colour = self._band_colour(item.get("score", 0), maximum)
            self._card(slide, x, yy, cw, ch, self.c["paper-200"])
            self._rule(slide, x, yy, Inches(0.09), colour, thickness=Pt(int(ch / 12700)))
            self._text(
                slide, item.get("id", ""), x + Inches(0.26), yy + Inches(0.1), Inches(0.5),
                Inches(0.3), size=13, bold=True, color=colour,
            )
            label = item.get("label", "")
            label_w = cw - Inches(1.85)
            self._text(
                slide, label, x + Inches(0.78), yy + Inches(0.1), label_w, ch - Inches(0.2),
                size=self._fit(label, label_w, ch - Inches(0.2), 12, min_pt=9.5, where="scores"),
                color=self.c["ink"], anchor=MSO_ANCHOR.MIDDLE,
            )
            self._text(
                slide, f"{item.get('score', 0)}/{maximum}", x + cw - Inches(0.95),
                yy + Inches(0.1), Inches(0.72), ch - Inches(0.2), size=17, bold=True,
                align=PP_ALIGN.RIGHT, color=colour, anchor=MSO_ANCHOR.MIDDLE,
            )
        if notes:
            ny = int(top + rows * (ch + Inches(0.18)) - Inches(0.1))
            self._text(
                slide, notes, grid_x, min(ny, int(floor - note_h)), grid_w, note_h,
                size=11, color=self.c["slate"], italic=True, spacing=1.3,
            )
        self._chrome(slide)

    def findings(self, s):
        slide = self._slide(self.c["paper"])
        top = self._heading(slide, s["title"], s.get("subtitle"))
        items = s.get("findings", [])
        gap = Inches(0.2)
        avail = SLIDE_H - top - Inches(1.15)
        h = int((avail - gap * max(len(items) - 1, 0)) / max(len(items), 1))
        for i, f in enumerate(items):
            y = int(top + i * (h + gap))
            self._card(slide, MARGIN, y, CONTENT_W, h, self.c["paper-200"])
            self._rule(slide, MARGIN, y, Inches(0.1), self.c["signal"], thickness=Pt(int(h / 12700)))
            self._text(
                slide, str(f.get("n", i + 1)), MARGIN + Inches(0.3), y + Inches(0.14),
                Inches(0.5), Inches(0.5), size=22, bold=True, color=self.c["signal-700"],
            )
            body_w = CONTENT_W - Inches(1.2)
            body_h = h - Inches(0.72)
            self._text(
                slide, f.get("heading", ""), MARGIN + Inches(0.85), y + Inches(0.16),
                body_w, Inches(0.36), size=16, bold=True,
            )
            self._text(
                slide, f.get("text", ""), MARGIN + Inches(0.85), y + Inches(0.58),
                body_w, body_h,
                size=self._fit(f.get("text", ""), body_w, body_h, 13, where="findings"),
                color=self.c["slate"],
            )
        self._chrome(slide)

    def table(self, s):
        slide = self._slide(self.c["paper"])
        top = self._heading(slide, s["title"], s.get("subtitle"))
        cols = s.get("columns", [])
        rows = s.get("rows", [])
        note_h = Inches(0.55) if s.get("note") else 0
        avail = (SLIDE_H - top - Inches(1.1) - note_h) / 914400

        # PowerPoint grows a row to fit its text, so a requested table height is a
        # minimum and not a maximum — an over-long cell pushes the last rows off
        # the slide. Column widths are therefore fixed first, then the type size
        # is stepped down until the *estimated* wrapped height fits the space.
        weights = s.get("widths") or [1] * len(cols)
        total = sum(weights) or 1
        col_w = [int(CONTENT_W * w / total) for w in weights]
        pad = Inches(0.22)

        def measure(size):
            line_h = size * 1.3 * LEADING / 72
            heights = []
            for row in [cols] + rows:
                lines = 1
                for c in range(len(cols)):
                    text = str(row[c]) if c < len(row) else ""
                    w_in = max((col_w[c] - pad) / 914400, 0.4)
                    cpl = max(int((w_in * 72) / (GLYPH * size)), 1)
                    lines = max(lines, math.ceil(len(text) / cpl))
                heights.append(lines * line_h + 0.16)
            return heights

        size = 12.0
        heights = measure(size)
        while sum(heights) > avail and size > 8.0:
            size -= 0.5
            heights = measure(size)
        if sum(heights) > avail:
            self.overflows.append(f"slide {self.n + 1} (table): too many rows or too much cell text")

        shape = slide.shapes.add_table(
            len(rows) + 1, len(cols), MARGIN, top, CONTENT_W, Inches(sum(heights))
        )
        tbl = shape.table
        tbl.first_row = True
        for c, width in enumerate(col_w):
            tbl.columns[c].width = width
        for r, h in enumerate(heights):
            tbl.rows[r].height = Inches(h)

        def style(cell, text, *, bold, colour, fill):
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.11)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            para = cell.text_frame.paragraphs[0]
            for run in para.runs:
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.name = self.font
                run.font.color.rgb = colour

        for c, name in enumerate(cols):
            style(tbl.cell(0, c), str(name), bold=True, colour=self.c["paper"], fill=self.c["ink"])
        for r, row in enumerate(rows, start=1):
            fill = self.c["paper"] if r % 2 else self.c["paper-200"]
            for c in range(len(cols)):
                text = str(row[c]) if c < len(row) else ""
                bold = c == 0 and s.get("boldFirstColumn", True)
                style(tbl.cell(r, c), text, bold=bold, colour=self.c["ink"], fill=fill)

        if s.get("note"):
            # Sits under the table, not pinned to the slide foot — a two-line note
            # on a pinned box runs into the footer chrome.
            note_y = int(top + Inches(sum(heights)) + Inches(0.14))
            box_h = max(int(SLIDE_H - Inches(1.05)) - note_y, int(Inches(0.3)))
            self._text(
                slide, s["note"], MARGIN, note_y, CONTENT_W, box_h,
                size=self._fit(s["note"], CONTENT_W, box_h, 11.5, where="table note"),
                color=self.c["slate"], italic=True,
            )
        self._chrome(slide)

    def questions(self, s):
        slide = self._slide(self.c["ink"])
        top = self._heading(slide, s["title"], s.get("subtitle"), on_dark=True)
        blocks = s.get("blocks", [])
        gap = Inches(0.2)
        h = int((SLIDE_H - top - Inches(1.15) - gap * max(len(blocks) - 1, 0)) / max(len(blocks), 1))
        for i, b in enumerate(blocks):
            y = int(top + i * (h + gap))
            self._card(slide, MARGIN, y, CONTENT_W, h, self.c["ink-700"])
            tag = self._pill(
                slide, MARGIN + Inches(0.3), y + Inches(0.16), Inches(0.46), Inches(0.46),
                self.c["signal"],
            )
            p = tag.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(b.get("tag", chr(65 + i)))
            run.font.size = Pt(15)
            run.font.bold = True
            run.font.name = self.font
            run.font.color.rgb = self.c["ink"]
            self._text(
                slide, b.get("heading", ""), MARGIN + Inches(1.0), y + Inches(0.13),
                CONTENT_W - Inches(1.4), Inches(0.36), size=17, bold=True, color=self.c["paper"],
            )
            if b.get("prompt"):
                body_w = CONTENT_W - Inches(1.4)
                body_h = h - Inches(0.52)
                self._text(
                    slide, b["prompt"], MARGIN + Inches(1.0), y + Inches(0.5), body_w, body_h,
                    size=self._fit(b["prompt"], body_w, body_h, 13, where="questions"),
                    color=RGBColor.from_string("A8BACD"), italic=True,
                )
        self._chrome(slide, on_dark=True)

    def closing(self, s):
        slide = self._slide(self.c["paper"])
        top = self._heading(slide, s["title"], s.get("subtitle"))
        items = s.get("items", [])
        gap = Inches(0.18)
        avail = SLIDE_H - top - Inches(1.3) - (Inches(0.34) if s.get("footerNote") else 0)
        h = int((avail - gap * max(len(items) - 1, 0)) / max(len(items), 1))
        for i, item in enumerate(items):
            y = int(top + i * (h + gap))
            self._card(slide, MARGIN, y, CONTENT_W, h, self.c["paper-200"])
            self._text(
                slide, item.get("lead", ""), MARGIN + Inches(0.32), y + Inches(0.14),
                Inches(3.4), h - Inches(0.28), size=14, bold=True, anchor=MSO_ANCHOR.MIDDLE,
            )
            body_w = CONTENT_W - Inches(4.3)
            body_h = h - Inches(0.28)
            self._text(
                slide, item.get("text", ""), MARGIN + Inches(3.9), y + Inches(0.14),
                body_w, body_h,
                size=self._fit(item.get("text", ""), body_w, body_h, 13, where="closing"),
                color=self.c["slate"], anchor=MSO_ANCHOR.MIDDLE,
            )
        if s.get("footerNote"):
            self._text(
                slide, s["footerNote"], MARGIN, SLIDE_H - Inches(1.32), CONTENT_W, Inches(0.5),
                size=12, color=self.c["ink"], bold=True,
            )
        self._chrome(slide)

    # ---------- driver ----------

    LAYOUTS = {
        "title": "title",
        "statement": "statement",
        "bullets": "bullets",
        "two-col": "two_col",
        "scores": "scores",
        "findings": "findings",
        "table": "table",
        "questions": "questions",
        "closing": "closing",
    }

    def build(self, out: Path) -> int:
        slides = self.spec.get("slides", [])
        if not slides:
            raise SystemExit("spec has no slides")
        for i, s in enumerate(slides, start=1):
            layout = s.get("layout")
            method = self.LAYOUTS.get(layout)
            if not method:
                raise SystemExit(
                    f"slide {i}: unknown layout {layout!r} — known: {', '.join(sorted(self.LAYOUTS))}"
                )
            getattr(self, method)(s)
        out.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(out))
        return len(slides)


def main() -> None:
    ap = argparse.ArgumentParser(description="Build an Evidura-branded deck from a JSON spec")
    ap.add_argument("spec", type=Path, help="deck spec JSON")
    ap.add_argument("-o", "--out", type=Path, required=True, help="output .pptx path")
    ap.add_argument("--font", default="Inter", help="typeface (brand default: Inter)")
    ap.add_argument("--tokens", type=Path, default=TOKENS, help="path to tokens.css")
    args = ap.parse_args()

    spec = json.loads(args.spec.read_text(encoding="utf-8"))
    deck = Deck(spec, args.font, load_tokens(args.tokens))
    count = deck.build(args.out)
    print(f"{args.out}  ({count} slides, {args.font}, tokens from {args.tokens.name})")
    for warning in deck.overflows:
        print(f"  overflow: {warning}", file=sys.stderr)
    if deck.overflows:
        sys.exit(2)


if __name__ == "__main__":
    main()
